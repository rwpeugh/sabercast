"""demo/generate_pitch_slide.py — produce a one-slide pitch deck for Sabercast.

Output: docs/demo/Sabercast_Pitch_Slide.pptx

The slide is the image-rich version of the pitch deck: title + hero stat band
+ product screenshot + concise tech column + vendor-logo footer. Layout
mirrors ``render_pitch_slide_png.py`` so the PPTX and PNG outputs match.

Editable in PowerPoint, Keynote, Google Slides, LibreOffice, etc.

To export a PNG for sharing, either re-run ``render_pitch_slide_png.py`` (the
canonical PNG), or open the PPTX and File -> Export -> PNG.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from PIL import Image


ROOT     = Path(__file__).resolve().parent.parent
OUT_PATH = ROOT / "docs" / "demo" / "Sabercast_Pitch_Slide.pptx"

LOGOS = ROOT / "docs" / "demo" / "logos"
SHOTS = ROOT / "docs" / "checkpoint3"

# Colour palette — match the PNG renderer
NAVY    = RGBColor(0x10, 0x2A, 0x55)
SLATE   = RGBColor(0x33, 0x44, 0x55)
GRAY    = RGBColor(0x66, 0x70, 0x80)
LIGHT   = RGBColor(0xF0, 0xF2, 0xF5)
CREAM   = RGBColor(0xFA, 0xFB, 0xFC)
ACCENT  = RGBColor(0x3A, 0x82, 0xCD)
GOOD    = RGBColor(0x2E, 0x86, 0x3E)
GOOD_BG = RGBColor(0xE8, 0xF4, 0xEA)
NEUTRAL = RGBColor(0xC0, 0x8C, 0x00)
NEUT_BG = RGBColor(0xFB, 0xF3, 0xDC)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)


# ─────────────────────────────────────────────────────────────────────────
# Drawing helpers
# ─────────────────────────────────────────────────────────────────────────

def _add_text(slide, x, y, w, h, text: str, *,
              size: float = 14, bold: bool = False, italic: bool = False,
              color: RGBColor = SLATE, align: PP_ALIGN = PP_ALIGN.LEFT,
              vcenter: bool = False, font: str = "Calibri") -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    if vcenter:
        from pptx.enum.text import MSO_ANCHOR
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color


def _add_box(slide, x, y, w, h, fill: RGBColor, *,
             line: RGBColor | None = None, lw: float = 0.5,
             rounded: bool = False) -> None:
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type,
                                    Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(lw)
    shape.shadow.inherit = False
    # Mute the default rounded-corner radius — python-pptx defaults to ~16% which
    # is far too rounded for a stat card.
    if rounded:
        try:
            shape.adjustments[0] = 0.10
        except Exception:
            pass


def _add_image(slide, path: Path, x: float, y: float, w: float, h: float,
               *, preserve_aspect: bool = True) -> None:
    if not path.exists():
        print(f"  warn: missing image {path}")
        return
    if preserve_aspect:
        with Image.open(path) as img:
            iw, ih = img.size
        img_aspect = iw / ih
        slot_aspect = w / h
        if img_aspect > slot_aspect:
            actual_w = w
            actual_h = w / img_aspect
        else:
            actual_h = h
            actual_w = h * img_aspect
        cx = x + (w - actual_w) / 2
        cy = y + (h - actual_h) / 2
        slide.shapes.add_picture(str(path), Inches(cx), Inches(cy),
                                 width=Inches(actual_w), height=Inches(actual_h))
    else:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y),
                                 width=Inches(w), height=Inches(h))


def _stat_card(slide, x: float, y: float, w: float, h: float,
               headline: str, label: str, detail: str, *,
               accent: RGBColor = GOOD, bg: RGBColor = GOOD_BG) -> None:
    """One big-number stat card — rounded panel + headline + sublabel."""
    _add_box(slide, x, y, w, h, fill=bg, line=accent, lw=0.75, rounded=True)
    # Big headline
    _add_text(slide, x, y + 0.10, w, 0.55, headline,
              size=28, bold=True, color=accent, align=PP_ALIGN.CENTER)
    # Label
    _add_text(slide, x, y + 0.66, w, 0.25, label,
              size=10, bold=True, color=SLATE, align=PP_ALIGN.CENTER)
    # Detail (p-value)
    _add_text(slide, x, y + 0.88, w, 0.22, detail,
              size=8.5, italic=True, color=GRAY, align=PP_ALIGN.CENTER)


def _logo_strip(slide, y: float, logo_h: float = 0.40) -> None:
    """Centered Powered-by row with 5 logos."""
    items = [
        ("openai.png.png",                  0.80),
        ("together_ai.png.png",             0.80),
        ("chromadb.png",                    0.65),
        ("streamlit.png.png",               0.65),
        ("Baseball_Reference_Logo.svg.png", 1.30),
    ]
    label_w = 1.15
    gap     = 0.40
    slide_w = 13.333
    total_w = label_w + sum(w for _, w in items) + gap * len(items)
    start_x = (slide_w - total_w) / 2

    _add_text(slide, start_x, y, label_w, logo_h,
              "Powered by", size=11, bold=True, color=NAVY,
              align=PP_ALIGN.CENTER, vcenter=True)

    x_cursor = start_x + label_w + gap
    for fname, w in items:
        _add_image(slide, LOGOS / fname, x_cursor, y, w, logo_h)
        x_cursor += w + gap


# ─────────────────────────────────────────────────────────────────────────
# Main slide builder
# ─────────────────────────────────────────────────────────────────────────

def build_slide() -> None:
    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # ─── BACKGROUND: cream wash (full bleed) ────────────────────────────
    _add_box(slide, 0, 0, 13.333, 7.5, fill=CREAM)

    # ─── BACKGROUND: subtle diamond watermark, upper-right ──────────────
    _add_image(slide, LOGOS / "baseball_diamond.png",
               x=10.0, y=0.15, w=3.10, h=3.10)

    # ─── LEFT EDGE: vertical baseball-seam strand ───────────────────────
    _add_image(slide, LOGOS / "seam_strand.png",
               x=0.10, y=0.20, w=0.40, h=7.10, preserve_aspect=False)

    # ─── TITLE BLOCK ────────────────────────────────────────────────────
    _add_text(slide, 0.75, 0.30, 9.0, 0.85, "Sabercast",
              size=42, bold=True, color=NAVY)
    _add_text(slide, 0.75, 1.10, 11.0, 0.35,
              "From roster gaps to ranked free-agent targets — in 13 seconds.",
              size=15, bold=True, color=SLATE)
    _add_text(slide, 0.75, 1.42, 11.0, 0.35,
              "Plus lineup planning and series scouting. Built for small / mid-market "
              "MLB front offices needing analyst leverage without a 20-person R&D shop.",
              size=11.5, italic=True, color=GRAY)
    _add_box(slide, 0.75, 1.78, 11.85, 0.025, ACCENT)

    # ─── HERO STAT BAND (4 cards) ───────────────────────────────────────
    cards_y = 1.95
    cards_h = 1.15
    band_x  = 0.75
    band_w  = 11.85
    gap     = 0.20
    card_w  = (band_w - 3 * gap) / 4
    cards = [
        ("13 sec", "Gap diagnosis to targets", "12 LLM calls parallelized",        GOOD, GOOD_BG),
        ("+70pp",  "RAG accuracy lift",        "20 Qs · McNemar p=0.0005",         GOOD, GOOD_BG),
        ("3.1×",   "Precision@10 lift",        "n=43 · p<0.0001 · finds signings", GOOD, GOOD_BG),
        ("59.9%",  "Position-gap hit rate",    "p=0.012 · 2B 74% (p=0.011)",       GOOD, GOOD_BG),
    ]
    for i, (head, label, detail, accent, bg) in enumerate(cards):
        x = band_x + i * (card_w + gap)
        _stat_card(slide, x, cards_y, card_w, cards_h, head, label, detail,
                   accent=accent, bg=bg)

    # ─── MAIN BODY: screenshot (LEFT) + tech column (RIGHT) ─────────────
    body_y    = 3.35
    body_h    = 2.45
    screen_x  = 0.75
    screen_w  = 7.45
    techcol_x = 8.45
    techcol_w = 4.15

    # Screenshot frame
    _add_box(slide, screen_x - 0.04, body_y - 0.04, screen_w + 0.08, body_h + 0.08,
             fill=WHITE, line=ACCENT, lw=0.75, rounded=True)
    _add_image(slide, SHOTS / "03_top_gap_card_with_candidates.png",
               x=screen_x, y=body_y, w=screen_w, h=body_h, preserve_aspect=True)
    _add_text(slide, screen_x, body_y + body_h + 0.05, screen_w, 0.22,
              "Position gap diagnosed  ·  3 free-agent targets ranked  ·  3 pricing comparables",
              size=9, italic=True, color=GRAY)

    # Right column — three concise blocks (Block 1 is outcome-framed for buyers)
    col_top = body_y + 0.00
    _add_text(slide, techcol_x, col_top + 0.00, techcol_w, 0.25,
              "THREE GM-OFFICE JOBS", size=10.5, bold=True, color=ACCENT)
    _add_text(slide, techcol_x, col_top + 0.27, techcol_w, 0.75,
              "• Plan tonight's lineup vs. the opponent\n"
              "• Scout the team you're facing\n"
              "• Diagnose roster gaps + rank FA targets",
              size=10, color=SLATE)

    _add_text(slide, techcol_x, col_top + 1.00, techcol_w, 0.25,
              "DATA + LLM STACK", size=10.5, bold=True, color=ACCENT)
    _add_text(slide, techcol_x, col_top + 1.27, techcol_w, 0.75,
              "pybaseball · Spotrac · Statcast\n"
              "gpt-4o + 4o-mini · text-embed-3\n"
              "ChromaDB RAG (999 players)",
              size=10, color=SLATE)

    _add_text(slide, techcol_x, col_top + 2.00, techcol_w, 0.25,
              "BUILD ECONOMICS", size=10.5, bold=True, color=ACCENT)
    _add_text(slide, techcol_x, col_top + 2.27, techcol_w, 0.30,
              "9-day build  ·  ~$47 platform spend",
              size=10, color=SLATE)

    # ─── HONEST-FRAMING LINE ───────────────────────────────────────────
    framing_y = 6.15
    _add_box(slide, 0.75, framing_y, 11.85, 0.025, NAVY)
    _add_text(slide, 0.75, framing_y + 0.10, 7.0, 0.30,
              "Decision-support tool — not a wins forecaster.",
              size=12, bold=True, color=NAVY)
    _add_text(slide, 7.0, framing_y + 0.12, 5.60, 0.30,
              "10 pre-registered tests  ·  4 significant  ·  5 honest nulls reported",
              size=10, italic=True, color=GRAY, align=PP_ALIGN.RIGHT)

    # ─── FOOTER: Powered-by logos + URL strip ──────────────────────────
    _logo_strip(slide, y=6.70, logo_h=0.40)

    _add_box(slide, 0.75, 7.15, 11.85, 0.015, GRAY)
    _add_text(slide, 0.75, 7.20, 5.5, 0.25,
              "Try the live app  ·  sabercast-mlb.streamlit.app",
              size=10, bold=True, color=NAVY)
    _add_text(slide, 5.5, 7.20, 2.5, 0.25,
              "MKTG 569  ·  Spring 2026",
              size=9, italic=True, color=GRAY, align=PP_ALIGN.CENTER)
    _add_text(slide, 7.5, 7.20, 5.10, 0.25,
              "github.com/rwpeugh/sabercast",
              size=10, bold=True, color=NAVY, align=PP_ALIGN.RIGHT)

    prs.save(str(OUT_PATH))
    print(f"saved {OUT_PATH}")
    print(f"size: {OUT_PATH.stat().st_size / 1024:.1f} KB")


if __name__ == "__main__":
    build_slide()
